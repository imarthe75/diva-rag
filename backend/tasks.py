# backend/tasks.py

import os
import json
import logging
from datetime import datetime
from uuid import UUID as UUIDType # Use UUIDType to avoid clash with uuid.uuid4
import uuid # For generating new UUIDs
from celery import Celery

# --- Gevent/Eventlet Monkey Patch (add this at the very top if using these pools) ---
# from gevent import monkey
# monkey.patch_all()
# If you are using --pool=prefork, you don't need the gevent monkey patch.
# If you plan to switch to gevent or eventlet, uncomment the two lines above
# and ensure 'gevent' or 'eventlet' is in your requirements.txt.

# --- SQLAlchemy and Models Imports ---
from database import get_db # Import the database session context manager
from models import DocumentVersion, DocumentChunk # Import your SQLAlchemy models

# --- External Libraries ---
from cryptography.fernet import Fernet
import requests
from minio import Minio
from minio.error import S3Error
from pypdf import PdfReader
import io
import subprocess
import tempfile
import pytesseract
from PIL import Image

# --- Logger Configuration ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# --- Celery Configuration ---
CELERY_BROKER_URL = os.getenv('CELERY_BROKER_URL', 'redis://localhost:6379/0')
CELERY_RESULT_BACKEND = os.getenv('CELERY_RESULT_BACKEND', 'redis://localhost:6379/0')
celery_app = Celery('tasks', broker=CELERY_BROKER_URL, backend=CELERY_RESULT_BACKEND)

# --- Environment Variables (Ensuring they are loaded correctly) ---
# These should ideally be loaded once at application startup or via your Docker setup.
# In a Celery worker, they are typically available via the Docker container's environment.
CEPH_ENDPOINT_URL = os.getenv('CEPH_ENDPOINT_URL')
CEPH_ACCESS_KEY = os.getenv('CEPH_ACCESS_KEY')
CEPH_SECRET_KEY = os.getenv('CEPH_SECRET_KEY')
CEPH_BUCKET_NAME = os.getenv('CEPH_BUCKET_NAME')

SYSTEM_MASTER_KEY = os.getenv('DOCUMENT_ENCRYPTION_KEY')
if not SYSTEM_MASTER_KEY:
    raise ValueError("DOCUMENT_ENCRYPTION_KEY is not configured in environment variables.")

OLLAMA_API_BASE_URL = os.getenv("OLLAMA_API_BASE_URL", "http://ollama:11434")
OLLAMA_EMBEDDING_MODEL = os.getenv("OLLAMA_EMBEDDING_MODEL", "nomic-embed-text")
OLLAMA_GENERATION_MODEL = os.getenv("OLLAMA_GENERATION_MODEL", "phi3:3.8b-mini-4k-instruct-q4_K_M")
OLLAMA_GENERATION_TIMEOUT = int(os.getenv("OLLAMA_GENERATION_TIMEOUT", "1200"))

# --- Utility Functions (consider moving these to a 'utils' directory) ---

# REMOVED: get_db_connection() - No longer needed with SQLAlchemy ORM

def get_s3_client():
    minio_host = CEPH_ENDPOINT_URL.replace("http://", "").replace("https://", "")
    secure_connection = CEPH_ENDPOINT_URL.startswith("https://")
    return Minio(
        minio_host,
        access_key=CEPH_ACCESS_KEY,
        secret_key=CEPH_SECRET_KEY,
        secure=secure_connection
    )

def encrypt(data: bytes, key: bytes) -> bytes:
    f = Fernet(key)
    return f.encrypt(data)

def decrypt(data: bytes, key: bytes) -> bytes:
    f = Fernet(key)
    return f.decrypt(data)

def get_ollama_embedding(text: str, model_name: str):
    headers = {'Content-Type': 'application/json'}
    data = {
        "model": model_name,
        "prompt": text
    }
    try:
        response = requests.post(f"{OLLAMA_API_BASE_URL}/api/embeddings", headers=headers, json=data, timeout=OLLAMA_GENERATION_TIMEOUT)
        response.raise_for_status()
        return response.json()['embedding']
    except requests.exceptions.Timeout as e:
        logger.error(f"Tiempo de espera agotado al obtener embedding de Ollama en {OLLAMA_API_BASE_URL}/api/embeddings: {e}")
        raise
    except requests.exceptions.RequestException as e:
        logger.error(f"Error al comunicarse con Ollama para embedding en {OLLAMA_API_BASE_URL}/api/embeddings: {e}")
        raise
    except Exception as e:
        logger.error(f"Error inesperado al obtener embedding de Ollama: {e}")
        raise

def get_ollama_generation(prompt: str, model_name: str):
    headers = {'Content-Type': 'application/json'}
    data = {
        "model": model_name,
        "prompt": prompt,
        "stream": False
    }
    try:
        logger.info(f"Solicitando generación para el modelo '{model_name}' (prompt: {prompt[:100]}...) con timeout {OLLAMA_GENERATION_TIMEOUT}s")
        response = requests.post(f"{OLLAMA_API_BASE_URL}/api/generate", headers=headers, json=data, timeout=OLLAMA_GENERATION_TIMEOUT)
        response.raise_for_status()
        return response.json()['response']
    except requests.exceptions.Timeout as e:
        logger.error(f"Tiempo de espera agotado al generar respuesta de Ollama en {OLLAMA_API_BASE_URL}/api/generate: {e}")
        raise
    except requests.exceptions.RequestException as e:
        logger.error(f"Error al comunicarse con Ollama en {OLLAMA_API_BASE_URL}/api/generate: {e}")
        raise
    except Exception as e:
        logger.error(f"Error inesperado al obtener generación de Ollama: {e}")
        raise

def extract_text_from_file_content(file_content_bytes: bytes, filename: str) -> str:
    _, file_extension = os.path.splitext(filename)
    file_extension = file_extension.lower()

    if file_extension == '.pdf':
        try:
            reader = PdfReader(io.BytesIO(file_content_bytes))
            text = ""
            for page in reader.pages:
                text += page.extract_text() or ""
            return text
        except Exception as e:
            logger.error(f"Error al extraer texto de PDF {filename}: {e}", exc_info=True)
            raise
    elif file_extension == '.txt':
        return file_content_bytes.decode('utf-8')
    elif file_extension == '.mobi':
        try:
            from mobi import Mobi # Requires 'mobi' package
            mobi_book = Mobi(io.BytesIO(file_content_bytes))
            mobi_book.parse()
            text_content = ""
            for chapter in mobi_book.contents:
                text_content += chapter.content.decode('utf-8', errors='ignore') + "\n"
            return text_content
        except ImportError:
            logger.error("La librería 'mobi' no está instalada. No se puede procesar .mobi")
            raise
        except Exception as e:
            logger.error(f"Error al extraer texto de MOBI {filename}: {e}", exc_info=True)
            raise
    elif file_extension == '.docx':
        try:
            from docx import Document # Requires 'python-docx' package
            document = Document(io.BytesIO(file_content_bytes))
            text = '\n'.join([paragraph.text for paragraph in document.paragraphs])
            return text
        except ImportError:
            logger.error("La librería 'python-docx' no está instalada. No se puede procesar .docx")
            raise
        except Exception as e:
            logger.error(f"Error al extraer texto de DOCX {filename}: {e}", exc_info=True)
            raise
    elif file_extension == '.xlsx':
        try:
            from openpyxl import load_workbook # Requires 'openpyxl' package
            workbook = load_workbook(io.BytesIO(file_content_bytes))
            text = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text.append(f"--- Hoja: {sheet_name} ---")
                for row in sheet.iter_rows():
                    row_values = [str(cell.value) if cell.value is not None else "" for cell in row]
                    text.append('\t'.join(row_values))
            return '\n'.join(text)
        except ImportError:
            logger.error("La librería 'openpyxl' no está instalada. No se puede procesar .xlsx")
            raise
        except Exception as e:
            logger.error(f"Error al extraer texto de XLSX {filename}: {e}", exc_info=True)
            raise
    elif file_extension == '.pptx':
        try:
            from pptx import Presentation # Requires 'python-pptx' package
            prs = Presentation(io.BytesIO(file_content_bytes))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return '\n'.join(text_runs)
        except ImportError:
            logger.error("La librería 'python-pptx' no está instalada. No se puede procesar .pptx")
            raise
        except Exception as e:
            logger.error(f"Error al extraer texto de PPTX {filename}: {e}", exc_info=True)
            raise
    elif file_extension == '.epub':
        try:
            from ebooklib import epub # Requires 'EbookLib' package
            import html2text # Requires 'html2text' package
            book = epub.read_epub(io.BytesIO(file_content_bytes))
            text_content = []
            for item in book.get_items():
                if item.get_type() == epub.ITEM_DOCUMENT:
                    text_content.append(html2text.html2text(item.get_content().decode('utf-8', errors='ignore')))
            return '\n'.join(text_content)
        except ImportError:
            logger.error("Las librerías 'EbookLib' o 'html2text' no están instaladas. No se puede procesar .epub")
            raise
        except Exception as e:
            logger.error(f"Error al extraer texto de EPUB {filename}: {e}", exc_info=True)
            raise
    elif file_extension == '.azw3':
        temp_input_azw3_path = None
        temp_output_txt_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.azw3', dir='/tmp') as temp_input:
                temp_input.write(file_content_bytes)
                temp_input_azw3_path = temp_input.name

            with tempfile.NamedTemporaryFile(delete=False, suffix='.txt', dir='/tmp') as temp_output:
                temp_output_txt_path = temp_output.name

            # Ensure calibre's ebook-convert is installed and in PATH
            command = ["ebook-convert", temp_input_azw3_path, temp_output_txt_path]
            result = subprocess.run(command, capture_output=True, text=True, check=True)

            if result.returncode != 0:
                logger.error(f"ebook-convert failed for {filename}. Stderr: {result.stderr}")
                raise Exception(f"ebook-convert failed: {result.stderr}")

            with open(temp_output_txt_path, 'r', encoding='utf-8') as f:
                text_content = f.read()
            return text_content

        except FileNotFoundError:
            logger.error("ebook-convert (Calibre) no encontrado. Asegúrate de que esté instalado en el contenedor.")
            raise
        except subprocess.CalledProcessError as e:
            logger.error(f"Error en ebook-convert para {filename}: {e}. Salida: {e.stdout}. Error: {e.stderr}", exc_info=True)
            raise
        except Exception as e:
            logger.error(f"Error al extraer texto de AZW3 {filename}: {e}", exc_info=True)
            raise
        finally:
            if temp_input_azw3_path and os.path.exists(temp_input_azw3_path):
                os.remove(temp_input_azw3_path)
            if temp_output_txt_path and os.path.exists(temp_output_txt_path):
                os.remove(temp_output_txt_path)
    elif file_extension in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff']:
        try:
            image = Image.open(io.BytesIO(file_content_bytes))
            # Ensure Tesseract is installed and available in PATH within the container
            text = pytesseract.image_to_string(image, lang='spa+eng')
            return text.strip()
        except pytesseract.TesseractNotFoundError:
            logger.error("Tesseract OCR no encontrado. Asegúrate de que esté instalado en el sistema y en el PATH.")
            raise
        except Exception as e:
            logger.error(f"Error al realizar OCR en la imagen {filename}: {e}", exc_info=True)
            raise
    else:
        logger.warning(f"Tipo de archivo no soportado para extracción de texto: {filename}")
        # Consider raising an error here if unsupported files should always fail RAG.
        # For now, returning empty string for unsupported types, but updating DB status is crucial.
        return ""

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 100) -> list[str]:
    chunks = []
    if not text:
        return chunks

    start_index = 0
    while start_index < len(text):
        end_index = min(start_index + chunk_size, len(text))
        chunks.append(text[start_index:end_index])
        start_index += (chunk_size - overlap)
        if start_index >= len(text) and start_index < len(text) + overlap and len(text[start_index:end_index].strip()) > 0:
            # Ensure the last chunk is not added if it's empty due to overlap
            if text[start_index:end_index].strip():
                chunks.append(text[start_index:end_index])
            break
        # Prevent infinite loop if chunk_size - overlap is 0 or negative
        if chunk_size - overlap <= 0:
            if start_index < len(text): # Add remaining text if any
                chunks.append(text[start_index:])
            break
    return chunks

# --- Celery Task for RAG Indexing ---

@celery_app.task(bind=True, max_retries=3, default_retry_delay=60)
def index_document_for_rag(self, document_version_id_str: str):
    """
    Celery task to extract text from a document version,
    generate embeddings, and index the chunks into the vector database for RAG.
    Assumes the file has already been uploaded to MinIO and scanned for viruses.
    """
    logger.info(f"RAG: Starting indexing for document_version_id: {document_version_id_str}")

    # **Use SQLAlchemy ORM with the get_db() context manager**
    with get_db() as db_session:
        try:
            minio_client = get_s3_client() # Minio client can be instantiated here
            fernet_master = Fernet(SYSTEM_MASTER_KEY.encode('utf-8'))

            # 1. Retrieve document version metadata using ORM
            document_version = db_session.query(DocumentVersion).filter_by(id=UUIDType(document_version_id_str)).first()

            if not document_version:
                logger.error(f"RAG: Document version not found in DB for indexing {document_version_id_str}.")
                raise ValueError("Document version record not found for RAG indexing.")

            # Update status to 'processing' (using ORM)
            document_version.processed_status = 'processing' # Use 'processed_status' from models.py
            document_version.last_processed_at = datetime.now() # Update timestamp
            db_session.add(document_version)
            db_session.commit() # Commit here to update status early
            db_session.refresh(document_version) # Refresh object after commit

            # 2. Download and decrypt the file
            try:
                response = minio_client.get_object(CEPH_BUCKET_NAME, document_version.ceph_path)
                encrypted_content = response.read()
                response.close()
                response.release_conn()
                logger.info(f"RAG: Encrypted file '{document_version.ceph_path}' downloaded for indexing.")
            except S3Error as e:
                logger.error(f"RAG: Minio (S3) error downloading '{document_version.ceph_path}': {e}", exc_info=True)
                document_version.processed_status = 'failed_download'
                document_version.last_processed_at = datetime.now()
                db_session.add(document_version)
                db_session.commit()
                raise # Re-raise for Celery retry
            except Exception as e:
                logger.error(f"RAG: Unexpected error downloading '{document_version.ceph_path}': {e}", exc_info=True)
                document_version.processed_status = 'failed_download'
                document_version.last_processed_at = datetime.now()
                db_session.add(document_version)
                db_session.commit()
                raise

            try:
                # Ensure decryption key is bytes. It comes from DB as LargeBinary, which is bytes in Python.
                file_encryption_key = fernet_master.decrypt(document_version.encryption_key_encrypted)
                decrypted_content = decrypt(encrypted_content, file_encryption_key)
                logger.info(f"RAG: File '{document_version.original_filename}' decrypted successfully.")
            except Exception as e:
                logger.error(f"RAG: Error decrypting file '{document_version.original_filename}': {e}", exc_info=True)
                document_version.processed_status = 'failed_decryption'
                document_version.last_processed_at = datetime.now()
                db_session.add(document_version)
                db_session.commit()
                raise

            # 3. Extract text from the file
            logger.info(f"RAG: Extracting text from {document_version.original_filename}...")
            text_content = extract_text_from_file_content(decrypted_content, document_version.original_filename)

            if not text_content or text_content.strip() == "":
                logger.warning(f"RAG: No significant text extracted from {document_version.original_filename}. Marking as 'no_text_extracted'.")
                document_version.processed_status = 'no_text_extracted'
                document_version.last_processed_at = datetime.now()
                db_session.add(document_version)
                db_session.commit()
                return # No text to index

            logger.info(f"RAG: Text extracted (first 200 chars): {text_content[:200]}...")

            # 4. Chunk the text and generate embeddings
            chunks = chunk_text(text_content)
            logger.info(f"RAG: Text divided into {len(chunks)} chunks.")

            # Delete existing chunks for this document version before inserting new ones (using ORM)
            # This ensures idempotency for re-indexing
            db_session.query(DocumentChunk).filter_by(document_version_id=document_version.id).delete()
            # No need for commit here if new chunks are inserted and committed in the same transaction

            # Save new chunks and embeddings (using ORM)
            for i, chunk_text_content in enumerate(chunks):
                embedding = get_ollama_embedding(chunk_text_content, model_name=OLLAMA_EMBEDDING_MODEL)
                new_chunk = DocumentChunk(
                    document_version_id=document_version.id,
                    chunk_text=chunk_text_content,
                    chunk_embedding=embedding, # SQLAlchemy/pgvector handles list conversion
                    chunk_order=i
                )
                db_session.add(new_chunk)
            db_session.commit() # Commit all new chunk insertions
            logger.info(f"RAG: Document version {document_version_id_str} successfully indexed in vector DB.")

            # 5. Update the final document version status to 'indexed'
            document_version.processed_status = 'indexed'
            document_version.last_processed_at = datetime.now()
            db_session.add(document_version)
            db_session.commit() # Final status commit

        except Exception as e:
            # IMPORTANT: Rollback the session in case of any error
            db_session.rollback()
            logger.error(f"Celery RAG Task: Indexing error for document_version_id {document_version_id_str}: {e}", exc_info=True)

            try:
                # Attempt to update status to 'failed_indexing' before retrying
                # This ensures the DB reflects failure status even for transient errors that lead to retry.
                # If the current session is broken, create a new one to update the status.
                with get_db() as update_db_session:
                    failed_document_version = update_db_session.query(DocumentVersion).filter_by(id=UUIDType(document_version_id_str)).first()
                    if failed_document_version:
                        failed_document_version.processed_status = 'failed_indexing'
                        failed_document_version.last_processed_at = datetime.now()
                        update_db_session.add(failed_document_version)
                        update_db_session.commit()
                self.retry(exc=e, countdown=self.default_retry_delay)
            except self.MaxRetriesExceededError:
                logger.error(f"Celery RAG Task: Max retries exceeded for indexing document_version_id {document_version_id_str}. Marking as critical failure.")
                # The status was already set to 'failed_indexing' in the previous block.
                # If for some reason that commit failed, we attempt it again here
                # just in case, but it should typically be persistent.
                pass # Status should already be updated
            except Exception as update_e:
                logger.error(f"Celery RAG Task: CRITICAL error updating failure status in DB for {document_version_id_str}: {update_e}")
        finally:
            # The session is automatically closed by the 'with get_db()' context manager.
            pass
